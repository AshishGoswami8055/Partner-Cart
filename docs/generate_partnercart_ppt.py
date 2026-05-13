#!/usr/bin/env python3
"""
PartnerCart — Final Year Project PowerPoint (no embedded images; placeholder slides).

Run:  python docs/generate_partnercart_ppt.py
Requires: pip install python-pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BASE = Path(__file__).resolve().parent
OUTPUT = BASE / "PartnerCart-Final-Year-Project.pptx"

ACCENT = RGBColor(0x1F, 0x4E, 0x79)  # blue similar to sample index
PLACEHOLDER_BG = RGBColor(0xF2, 0xF2, 0xF2)


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Title
    box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1.2))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.bold = True
    r.font.size = Pt(40)
    r.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER
    # Subtitle
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1.5))
    tf2 = box2.text_frame
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = subtitle
    r2.font.size = Pt(20)
    p2.alignment = PP_ALIGN.CENTER


def add_section_slide(prs: Presentation, heading: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tit = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.7))
    tf = tit.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = heading
    r.font.bold = True
    r.font.size = Pt(28)
    r.font.color.rgb = ACCENT
    body = slide.shapes.add_textbox(Inches(0.6), Inches(1.15), Inches(8.8), Inches(5.5))
    btf = body.text_frame
    btf.margin_left = Inches(0.05)
    for i, line in enumerate(bullets):
        para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        para.text = line
        para.level = 0
        para.font.size = Pt(17)
        para.space_after = Pt(8)


def add_index_table_slide(prs: Presentation, rows: list[tuple[str, str, str]], title: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tit = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.55))
    tf = tit.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.bold = True
    r.font.size = Pt(30)
    r.font.color.rgb = ACCENT
    r.font.underline = True
    p.alignment = PP_ALIGN.CENTER

    num_rows = len(rows) + 1
    num_cols = 3
    left, top, width, height = Inches(0.45), Inches(0.95), Inches(9.1), Inches(5.6)
    table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table
    table.columns[0].width = Inches(0.95)
    table.columns[1].width = Inches(6.65)
    table.columns[2].width = Inches(1.5)

    headers = ("Sr. No.", "Topic", "Slide No.")
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.bold = True
                run.font.size = Pt(11)
                run.font.color.rgb = ACCENT
            para.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xDA, 0xEE, 0xF3)

    for i, (sr, topic, slide_no) in enumerate(rows, start=1):
        for j, val in enumerate((sr, topic, slide_no)):
            cell = table.cell(i, j)
            cell.text = val
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(10)
                if j == 0 and sr and sr[0].isdigit() and "." not in sr[:3]:
                    for run in para.runs:
                        run.font.bold = True
                        run.font.color.rgb = ACCENT
                para.alignment = PP_ALIGN.LEFT if j == 1 else PP_ALIGN.CENTER


def add_image_placeholder_slide(prs: Presentation, title: str, instructions: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tit = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.6))
    tf = tit.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.bold = True
    r.font.size = Pt(26)
    r.font.color.rgb = RGBColor(0xC0, 0x00, 0x00)
    box = slide.shapes.add_shape(1, Inches(0.7), Inches(1.1), Inches(8.6), Inches(4.9))  # rectangle
    box.fill.solid()
    box.fill.fore_color.rgb = PLACEHOLDER_BG
    box.line.color.rgb = ACCENT
    box.line.width = Pt(2)
    inner = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8.2), Inches(3))
    inf = inner.text_frame
    inf.word_wrap = True
    ip = inf.paragraphs[0]
    ir = ip.add_run()
    ir.text = "IMAGE PLACEHOLDER\n\n" + instructions
    ir.font.size = Pt(18)
    ir.font.italic = True
    ip.alignment = PP_ALIGN.CENTER


def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "PartnerCart",
        "Full-Stack Multi-Vendor Marketplace\nFinal Year Project Presentation\n"
        "(React · Vite · Node.js · Express · MongoDB · Razorpay)",
    )

    # Index part 1 — matches sample structure; PartnerCart wording; Slide No. = this deck
    index_part1: list[tuple[str, str, str]] = [
        ("01", "Introduction", "4"),
        ("1.1", "Project Overview", "4"),
        ("1.2", "Purpose of the System", "5"),
        ("1.3", "Scope of the Project", "6"),
        ("02", "Proposed System", "7"),
        ("2.1", "System Objectives", "7"),
        ("2.2", "Advantages", "8"),
        ("2.3", "Feasibility Study", "9"),
        ("2.3.1", "Technical Feasibility", "9"),
        ("2.3.2", "Economical Feasibility", "9"),
        ("2.3.3", "Operational Feasibility", "9"),
        ("03", "System Analysis", "10"),
        ("3.1", "Existing System", "10"),
        ("3.2", "Need for New System", "10"),
        ("3.3", "SRS (Software Requirement Specification)", "11"),
        ("04", "System Planning", "12"),
        ("4.1", "Requirement Analysis & Data Gathering", "12"),
        ("05", "Tools & Environment Used", "13"),
        ("5.1", "Hardware & Software Requirements", "13"),
        ("5.1.1", "Software Requirement (React, Vite, Node, Express, MongoDB, …)", "13"),
        ("5.1.2", "Hardware Requirement", "13"),
        ("5.2", "Development Tools (VS Code, Browser, Git, Postman)", "14"),
        ("06", "System Design", "15"),
        ("6.1", "Data Flow Diagram (DFD)", "15"),
        ("6.1.1", "Level 0 DFD (Context Diagram)", "16"),
        ("6.1.2", "Level 1 DFD (Customer & Vendor Processes)", "17"),
        ("6.1.3", "Level 2 DFD — Admin / Orders & Payments", "18"),
    ]

    index_part2: list[tuple[str, str, str]] = [
        ("6.2", "Database Design", "19"),
        ("6.2.1", "Tables Description (User, Vendor, Category, Product, Order, …)", "19"),
        ("", "ER Diagram (recommended supporting figure)", "20"),
        ("6.3", "System Flow (Process Explanation — Browse to Delivery)", "21"),
        ("07", "Modules Description", "22"),
        ("7.1", "Authentication & Role Management Module", "22"),
        ("7.2", "Admin Login & Governance Module", "23"),
        ("7.3", "Category & Catalog Management Module", "24"),
        ("7.4", "Order & Payment Module (COD / Razorpay)", "25"),
        ("08", "Implementation", "26"),
        ("8.1", "Project Workflow", "26"),
        ("8.2", "Code Logic Overview", "27"),
        ("09", "Output Screens (Results)", "28"),
        ("9.1", "Login / Sign-up Page", "28"),
        ("9.2", "Dashboard (Customer / Vendor / Admin)", "29"),
        ("9.3", "Marketplace / Category / Product Listing", "30"),
        ("9.4", "Checkout / Order / Payment Screen", "31"),
        ("10", "Advantages of the System", "32"),
        ("11", "Limitations", "33"),
        ("12", "Future Enhancements", "34"),
        ("13", "Conclusion", "35"),
        ("14", "References", "36"),
    ]

    add_index_table_slide(prs, index_part1, "INDEX (1 / 2)")
    add_index_table_slide(prs, index_part2, "INDEX (2 / 2)")

    add_section_slide(
        prs,
        "01. Introduction — 1.1 Project Overview",
        [
            "PartnerCart is an online multi-vendor marketplace connecting customers, sellers (vendors), and a platform admin.",
            "Customers browse products from many vendors in one cart, pay online (Razorpay) or choose COD, and track orders.",
            "Vendors manage store profile, catalogue, inventory, orders, coupons, and messaging.",
            "Admins onboard vendors, manage categories/users/products at scale, view analytics and audit trails.",
            "Technology: SPA built with React + Vite; REST API with Node.js & Express; data in MongoDB; realtime via Socket.IO.",
        ],
    )

    add_section_slide(
        prs,
        "01. Introduction — 1.2 Purpose of the System",
        [
            "Digitize buying and selling for small businesses under one trusted platform.",
            "Reduce fragmented communication (calls/WhatsApp) with structured orders, statuses, and notifications.",
            "Provide secure login, role-based dashboards, and auditable administrative actions.",
            "Demonstrate integration of modern payments (Razorpay) with a realistic e-commerce lifecycle.",
        ],
    )

    add_section_slide(
        prs,
        "01. Introduction — 1.3 Scope of the Project",
        [
            "In scope: Registration/login (customer, vendor applicant, dedicated admin portal), catalogue & search, cart & wishlist, checkout, COD/Razorpay, order lifecycle, reviews, coupons, messaging, basic analytics.",
            "Actors: Customer, Vendor, Administrator; external services: payment gateway, optional email (SMTP), optional media (Cloudinary).",
            "Out of scope (typical FY project boundary): Own payment gateway, native mobile apps, full logistics/carrier integrations, GST invoicing automation.",
        ],
    )

    add_section_slide(
        prs,
        "02. Proposed System — 2.1 System Objectives",
        [
            "Single marketplace UX with multi-vendor cart split into vendor-wise order groups at checkout.",
            "Reliable authentication: JWT access/refresh patterns, optional Google OAuth, OTP-based password flows.",
            "Operational clarity: dashboards for each role; admin moderation and vendor application workflow.",
            "Extensible persistence in MongoDB (products, variants, orders, payments, chats, audits).",
        ],
    )

    add_section_slide(
        prs,
        "02. Proposed System — 2.2 Advantages",
        [
            "24×7 storefronts without physical co-location of inventory.",
            "Centralized taxonomy and governance with vendor self-service catalogue tools.",
            "Multiple payment options improve conversion; COD supports cash-preferred users.",
            "Reusable SPA + API separation eases scaling and deployment.",
        ],
    )

    add_section_slide(
        prs,
        "02. Proposed System — 2.3 Feasibility Study",
        [
            "2.3.1 Technical: MERN-style stack widely documented; clear separation of frontend/backend; Razorpay & OAuth are standard integrations.",
            "2.3.2 Economical: Uses open-source runtimes; variable cost limited to hosting, DB tier, SMS/email sends, payment fees.",
            "2.3.3 Operational: Web browsers + moderate PC/laptop suffice; admin/vendor training via guided UIs.",
        ],
    )

    add_section_slide(
        prs,
        "03. System Analysis",
        [
            "3.1 Existing system: Informal ordering (phone/DM), isolated vendor sites, spreadsheets for stock — prone to errors, no unified trust.",
            "3.2 Need: One platform with identity, catalogue, carts, structured payments, order history, messaging, moderation.",
            "3.3 SRS (high level): Functional needs — auth CRUD per role; product CRUD; cart; checkout; payments; notifications; messaging; admin APIs.",
            "Non-functional: Security (HTTPS, hashing, sanitization), performance (indexed queries), availability aligned with hosting SLA.",
        ],
    )

    add_section_slide(
        prs,
        "03.3 SRS — Key Requirements (abbreviated)",
        [
            "FR: User registration/login; vendor application & admin approval; category CRUD (admin); product CRUD/bulk upload (vendor).",
            "FR: Cart & wishlist; place order COD or Razorpay; verify payment signatures; vendor order status updates.",
            "FR: Reviews; coupons apply; messaging threads; notifications; analytics endpoints per role.",
            "NFR: Role-based routing; rate limits on sensitive routes; input validation with shared schema discipline (e.g., Zod).",
        ],
    )

    add_section_slide(
        prs,
        "04. System Planning — 4.1 Requirement Analysis & Data Gathering",
        [
            "Studied stakeholder needs: browse → cart → checkout → pay → fulfil → notify.",
            "Derived entities: Users, Vendors, Categories, Products (variants), Cart lines, Orders (orderGroups per vendor), Payments, Coupons, Conversations.",
            "API surface grouped under `/api/v1`: auth, users, vendors, categories, products, cart, orders, wishlist, reviews, coupons, messages, notifications, admin, analytics.",
            "Environment configuration for client origin, secrets, Razorpay keys, SMTP, Cloudinary documented for deployment.",
        ],
    )

    add_section_slide(
        prs,
        "05. Tools & Environment — 5.1 Hardware & Software",
        [
            "5.1.1 Software: Client — React 18, Vite, Tailwind, Redux Toolkit, React Query, React Router, Axios, Socket.IO client.",
            "Server — Node.js 18+, Express, Mongoose, Passport (Google OAuth), Razorpay SDK, JWT, Multer, Nodemailer, Helmet/CORS/compression.",
            "Database: MongoDB (Atlas or local). Optional: Cloudinary for images.",
            "5.1.2 Hardware: Developer PC (8 GB+ RAM recommended), stable internet; standard x64 OS; modern browser for testing.",
        ],
    )

    add_section_slide(
        prs,
        "05.2 Development Tools",
        [
            "VS Code (or similar IDE), Git version control, browser DevTools, Postman/Thunder Client for API testing.",
            "Node package managers (npm); ESLint on server; Vite dev server + API proxy pattern for local integration.",
        ],
    )

    add_section_slide(
        prs,
        "06. System Design — 6.1 Data Flow Diagram (concept)",
        [
            "Context: Customers, Vendors, Admin, Razorpay, SMTP/optional storage interact with PartnerCart platform.",
            "Level 1: Authentication, catalogue, cart, checkout/payment, fulfilment/notifications, admin governance.",
            "Level 2: Order placement branches COD vs Razorpay; webhook/signature verification; status updates fan-out per vendor group.",
            "Formal diagrams drawn in draw.io / Lucidchart / Visio — see following placeholder slides.",
        ],
    )

    add_image_placeholder_slide(
        prs,
        "[Insert Figure] 6.1.1 — Level 0 DFD",
        "Paste CONTEXT diagram:\nOuter entities: Customer, Vendor, Admin, Razorpay, Email/SMTP, Cloudinary.\nCentral process: PartnerCart Platform.\nLabel major data flows (catalog browse, orders, payments, notifications).",
    )
    add_image_placeholder_slide(
        prs,
        "[Insert Figure] 6.1.2 — Level 1 DFD",
        "Show decomposition: Login/Register; Browse & Search; Cart; Checkout & Payment Adapter; Vendor Store Management; Admin Moderation;\nPlus data stores Users, Catalogue, Orders, Notifications.",
    )
    add_image_placeholder_slide(
        prs,
        "[Insert Figure] 6.1.3 — Level 2 DFD (Admin / Orders & Payments)",
        "Zoom into admin actions (approve vendor, manage categories) OR order/payment pipeline\n(validate cart → create order → Razorpay order → verify signature → mark paid → notify).",
    )

    add_section_slide(
        prs,
        "06.2 Database Design — 6.2.1 Tables (Collections) Overview",
        [
            "User — credentials, roles, profile, preferences, OAuth ids, security metadata.",
            "Vendor & VendorApplication — store details, application status, linkage to admin decisions.",
            "Category — hierarchy/slug for navigation; Product — vendor ref, variants, stock, images, text index for search.",
            "Cart / Wishlist — line items referencing products/variants; Order — orderNumber, embedded shippingAddress, orderGroups[].",
            "Also: Coupon, Review, Conversation/Message, Notification, Otp, AuditLog, AnalyticsSnapshot as supporting collections.",
        ],
    )

    add_image_placeholder_slide(
        prs,
        "[Insert Figure] ER Diagram (optional but recommended)",
        "Draw rectangles for each collection above and relationships (Vendor 1→N Product, Customer 1→N Order, Order N→embedded orderGroups keyed to Vendor).\nHighlight payment sub-document on Order.",
    )

    add_section_slide(
        prs,
        "06.3 System Flow — Process Explanation",
        [
            "1) Visitor browses marketplace / PDP → adds items to authenticated cart.",
            "2) Checkout validates address → applies coupon optionally → computes totals per vendor group.",
            "3a) COD: order confirmed unpaid path; 3b) Online: Razorpay order created → client completes payment → server verifies signature → paid.",
            "4) Vendors progress order statuses; customers see timeline; chats/notifications augment awareness.",
            "5) Admin oversees exceptions via dashboards and audit-friendly actions.",
        ],
    )

    add_section_slide(
        prs,
        "07. Modules — 7.1 Authentication & Role Management",
        [
            "Customer & vendor signup/login; JWT access + refresh; Google OAuth bridge; OTP forgot/reset and change-password via email.",
            "Frontend route guards: `/app` customer shell, `/vendor` vendor-only, `/admin` admin-only + separate admin login view.",
        ],
    )

    add_section_slide(
        prs,
        "07. Modules — 7.2 Admin Login & Governance",
        [
            "Admin authentication surface; dashboards for approvals (vendors/applications); user/product/category oversight.",
            "Audit logs and analytics APIs for accountable operations.",
        ],
    )

    add_section_slide(
        prs,
        "07. Modules — 7.3 Category & Catalog Management",
        [
            "Admin-maintained taxonomy; vendors create/update listings, variants, stock, imagery metadata; marketplace search/filter.",
            "Wishlist persistence; bulk product tooling where implemented for sellers.",
        ],
    )

    add_section_slide(
        prs,
        "07. Modules — 7.4 Order & Payment Module",
        [
            "Multi-vendor decomposition at order creation (per-vendor totals and statuses); customer order history.",
            "Razorpay prepare/verify/abandon lifecycle; COD path; cancellations where supported by business rules.",
        ],
    )

    add_section_slide(
        prs,
        "08. Implementation — 8.1 Project Workflow",
        [
            "Repository split: client (Vite React SPA) + server (Express mono-service API). Shared validation mindset (e.g., Zod on server routes).",
            "Local development: concurrent dev servers; `.env` for secrets; seed script for realistic demo catalogue.",
            "Production mental model: static client build served via CDN/nginx; API behind HTTPS with CORS allow-list.",
        ],
    )

    add_section_slide(
        prs,
        "08.2 Code Logic Overview",
        [
            "`createApp()` wires security middleware → `/api/v1` routers → centralized error handlers.",
            "Controllers validate input → services encapsulate Mongoose mutations (orders, inventory, payouts fields as designed).",
            "Client uses Redux Toolkit + React Query for server state/cache; Axios interceptors attach tokens.",
        ],
    )

    add_image_placeholder_slide(
        prs,
        "[Insert Figure] 9.1 — Login / Sign-up",
        "Screenshot: Unified login/signup (+ optional Forgot password OTP). Optional second shot: Dedicated Admin login page.",
    )
    add_image_placeholder_slide(
        prs,
        "[Insert Figure] 9.2 — Dashboard",
        "One slide each or collage: Customer home; Vendor KPI dashboard; Admin stats — pick clearest KPI cards/charts.",
    )
    add_image_placeholder_slide(
        prs,
        "[Insert Figure] 9.3 — Marketplace / Category / Products",
        "Screenshot: Marketplace listing with filters; category navigation; detailed product page with variants/add-to-cart.",
    )
    add_image_placeholder_slide(
        prs,
        "[Insert Figure] 9.4 — Checkout / Orders / Payment",
        "Screenshot: Checkout address + payment method; Razorpay modal or success; customer order list/vendor order queue.",
    )

    add_section_slide(
        prs,
        "10. Advantages of the System",
        [
            "Unified marketplace improves discovery; vendors retain separate fulfilment pipelines.",
            "Dual payments improve reach; realtime + email-style notifications deepen engagement.",
            "Role separation and audit hooks support safer operations for academic/demo production readiness.",
        ],
    )

    add_section_slide(
        prs,
        "11. Limitations",
        [
            "No native mobile apps in baseline web project; courier tracking integration not end-to-end.",
            "Operational concerns (scaling, redundancy) depend on chosen hosting—not fully modeled here.",
            "Advanced recommendation/fraud ML not included by default.",
        ],
    )

    add_section_slide(
        prs,
        "12. Future Enhancements",
        [
            "Carrier APIs, richer GST invoicing, mobile clients, multilingual UI, loyalty points, subscriptions, vector/semantic search.",
        ],
    )

    add_section_slide(
        prs,
        "13. Conclusion",
        [
            "PartnerCart demonstrates a complete mini-marketplace lifecycle with realistic roles, payments integration, and extensible persistence.",
            "The stack mirrors industry practice—react SPA + Node API + document DB—and is suitable as a Capstone FY project artefact.",
        ],
    )

    add_section_slide(
        prs,
        "14. References",
        [
            "Official docs: react.dev · vitejs.dev · expressjs.com · mongoosejs.com · mongodb.com/docs · razorpay.com/docs",
            "Standards/best practices: OWASP Cheat Sheets (JWT/CORS/auth), Razorpay server integration guides.",
            "Tutorial & learning resources as cited by your institute / guide.",
        ],
    )

    prs.save(str(OUTPUT))
    return OUTPUT


if __name__ == "__main__":
    path = build_presentation()
    print(f"Wrote: {path}")
